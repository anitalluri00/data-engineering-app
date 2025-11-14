import os
import uuid
import pandas as pd
import PyPDF2
from docx import Document
import pptx
import openpyxl
import csv
from PIL import Image
import speech_recognition as sr
from moviepy.editor import VideoFileClip
import zipfile
import logging
from database.models import DataStorage
from database.connection import db_connection
import mysql.connector

class FileProcessor:
    def __init__(self):
        self.db = DataStorage()
        self.supported_formats = {
            'text': ['.txt', '.rtf'],
            'documents': ['.pdf', '.doc', '.docx'],
            'spreadsheets': ['.xlsx', '.xls', '.csv'],
            'presentations': ['.ppt', '.pptx'],
            'images': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tif', '.tiff'],
            'audio': ['.mp3', '.wav', '.wma', '.m4a'],
            'video': ['.mp4', '.mov', '.wmv', '.flv', '.avi'],
            'archives': ['.zip', '.rar', '.7z'],
            'executables': ['.exe', '.dll', '.bat', '.sys']
        }
    
    def process_file(self, file_data, filename, source_type="upload"):
        """Process any file type and store in database"""
        try:
            file_id = str(uuid.uuid4())
            file_ext = os.path.splitext(filename)[1].lower()
            
            # Store file metadata and content
            metadata = {
                'file_extension': file_ext,
                'file_size': len(file_data),
                'source_type': source_type,
                'file_type': self._get_file_type(file_ext)
            }
            
            # Store file in database
            self._store_file(file_id, filename, file_ext, file_data, metadata)
            
            # Process file content based on type
            processed_data = self._extract_content(file_id, file_data, file_ext)
            
            return file_id, processed_data
            
        except Exception as e:
            logging.error(f"Error processing file {filename}: {e}")
            raise
    
    def _get_file_type(self, file_ext):
        for file_type, extensions in self.supported_formats.items():
            if file_ext in extensions:
                return file_type
        return 'unknown'
    
    def _store_file(self, file_id, filename, file_ext, file_data, metadata):
        try:
            conn = db_connection.get_connection()
            cursor = conn.cursor()
            
            cursor.execute("""
                INSERT INTO files (file_id, filename, file_type, file_size, source_type, content, metadata)
                VALUES (%s, %s, %s, %s, %s, %s, %s)
            """, (file_id, filename, file_ext, len(file_data), metadata['source_type'], file_data, str(metadata)))
            
            conn.commit()
            cursor.close()
            
        except mysql.connector.Error as e:
            logging.error(f"Error storing file in database: {e}")
            raise
    
    def _extract_content(self, file_id, file_data, file_ext):
        """Extract content based on file type"""
        try:
            content_type = self._get_file_type(file_ext)
            extracted_text = ""
            
            if content_type == 'text':
                extracted_text = file_data.decode('utf-8')
                
            elif content_type == 'documents':
                if file_ext == '.pdf':
                    extracted_text = self._extract_pdf_text(file_data)
                elif file_ext in ['.doc', '.docx']:
                    extracted_text = self._extract_doc_text(file_data)
                    
            elif content_type == 'spreadsheets':
                extracted_text = self._extract_spreadsheet_text(file_data, file_ext)
                
            elif content_type == 'presentations':
                extracted_text = self._extract_presentation_text(file_data)
                
            elif content_type == 'images':
                extracted_text = self._extract_image_metadata(file_data)
                
            elif content_type == 'audio':
                extracted_text = self._extract_audio_text(file_data)
                
            elif content_type == 'video':
                extracted_text = self._extract_video_metadata(file_data)
                
            # Store processed data
            data_id = str(uuid.uuid4())
            self._store_processed_data(data_id, file_id, content_type, extracted_text)
            
            return {
                'data_id': data_id,
                'content_type': content_type,
                'extracted_text': extracted_text[:1000] + "..." if len(extracted_text) > 1000 else extracted_text,
                'word_count': len(extracted_text.split()),
                'char_count': len(extracted_text)
            }
            
        except Exception as e:
            logging.error(f"Error extracting content: {e}")
            return {"error": str(e)}
    
    def _extract_pdf_text(self, file_data):
        import io
        pdf_file = io.BytesIO(file_data)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
        return text
    
    def _extract_doc_text(self, file_data):
        import io
        doc_file = io.BytesIO(file_data)
        doc = Document(doc_file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    
    def _extract_spreadsheet_text(self, file_data, file_ext):
        import io
        excel_file = io.BytesIO(file_data)
        
        if file_ext == '.csv':
            df = pd.read_csv(excel_file)
        else:
            df = pd.read_excel(excel_file)
        
        return df.to_string()
    
    def _extract_presentation_text(self, file_data):
        import io
        ppt_file = io.BytesIO(file_data)
        presentation = pptx.Presentation(ppt_file)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
    def _extract_image_metadata(self, file_data):
        import io
        image_file = io.BytesIO(file_data)
        image = Image.open(image_file)
        return f"Image: {image.format}, Size: {image.size}, Mode: {image.mode}"
    
    def _extract_audio_text(self, file_data):
        # Basic audio metadata extraction
        return "Audio content - transcription would require additional processing"
    
    def _extract_video_metadata(self, file_data):
        # Basic video metadata extraction
        return "Video content - analysis would require additional processing"
    
    def _store_processed_data(self, data_id, file_id, content_type, extracted_text):
        try:
            conn = db_connection.get_connection()
            cursor = conn.cursor()
            
            cursor.execute("""
                INSERT INTO processed_data (data_id, file_id, content_type, extracted_text, word_count, char_count)
                VALUES (%s, %s, %s, %s, %s, %s)
            """, (data_id, file_id, content_type, extracted_text, len(extracted_text.split()), len(extracted_text)))
            
            # Mark file as processed
            cursor.execute("""
                UPDATE files SET processed = TRUE WHERE file_id = %s
            """, (file_id,))
            
            conn.commit()
            cursor.close()
            
        except mysql.connector.Error as e:
            logging.error(f"Error storing processed data: {e}")
            raise